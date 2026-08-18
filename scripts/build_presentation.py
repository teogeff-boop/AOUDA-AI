"""
JARVIS / AOUDA — PowerPoint Generator (exact replica of ÖWF AMADEE-27 style)
=============================================================================
Replicates the exact design aesthetic of the ÖWF 'The Caillou Project' deck:
  - Clean white backgrounds
  - Monospace / Typewriter header styling (Consolas / Courier New)
  - Top-left & Top-right logo placement (ESTACA & ÖWF)
  - Full-bleed photo title slide with white banner overlay
  - Vertical step-by-step process sidebars
  - Generous whitespace, professional engineering typography
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ── Colors & Fonts ────────────────────────────────────────────────────────────
COLOR_DARK      = RGBColor(30, 30, 30)      # #1E1E1E
COLOR_MUTED     = RGBColor(100, 100, 100)  # #646464
COLOR_WHITE     = RGBColor(255, 255, 255)
COLOR_BG_CARD   = RGBColor(250, 248, 242)  # Light warm beige (like PDF sidebar)
COLOR_BORDER    = RGBColor(220, 215, 200)  # Subtle border
COLOR_ORANGE    = RGBColor(220, 120, 50)   # Numbered circle orange accent
COLOR_BLUE_ACC  = RGBColor(0, 102, 204)   # Blue accent
COLOR_RED_ACC   = RGBColor(196, 30, 58)    # Red accent

FONT_HEADING = "Consolas"
FONT_BODY    = "Calibri"

OWF_PATH    = Path("data/owf_logo.png")
ESTACA_PATH = Path("data/estaca_logo.png")
BG_TITLE    = Path("data/bg_title.jpg")

# ── Utility Helpers ───────────────────────────────────────────────────────────

def add_header(slide, title_text):
    """Standard header matching PDF slides: logos at top left & right, centered title."""
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_WHITE
    bg.line.fill.background()

    # Left Logo (ESTACA)
    if ESTACA_PATH.exists():
        slide.shapes.add_picture(str(ESTACA_PATH), Inches(0.4), Inches(0.25), height=Inches(0.85))

    # Right Logo (ÖWF)
    if OWF_PATH.exists():
        slide.shapes.add_picture(str(OWF_PATH), Inches(11.6), Inches(0.25), height=Inches(0.85))

    # Centered Header Title
    tb = slide.shapes.add_textbox(Inches(2.2), Inches(0.35), Inches(8.933), Inches(0.75))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_HEADING
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK


def add_text_paragraph(slide, text, l, t, w, h, font_name=FONT_BODY, font_size=15,
                       bold=False, color=COLOR_DARK, align=PP_ALIGN.LEFT, space_after=Pt(14)):
    """Add a text block with paragraph formatting."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_after = space_after
    return tb


def add_bullet_list(slide, items, l, t, w, h, font_size=14, space_after=Pt(16)):
    """Add bullet points with bold lead-ins (e.g. 'Lead-in: details')."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True

    for idx, (title, detail) in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = space_after
        p.level = 0
        
        # Bullet symbol & title
        r1 = p.add_run()
        r1.text = f"•  {title}: "
        r1.font.name = FONT_BODY
        r1.font.size = Pt(font_size)
        r1.font.bold = True
        r1.font.color.rgb = COLOR_DARK

        # Detail text
        r2 = p.add_run()
        r2.text = detail
        r2.font.name = FONT_BODY
        r2.font.size = Pt(font_size)
        r2.font.bold = False
        r2.font.color.rgb = COLOR_DARK


def add_process_sidebar(slide, steps, l=8.8, t=1.35, w=4.0, h=5.7):
    """Draw vertical 6-step workflow column (like right side of slides 4-6 in reference PDF)."""
    # Card Background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_BG_CARD
    card.line.color.rgb = COLOR_BORDER
    card.line.width = Pt(1)

    step_h = h / len(steps)
    for i, (num_str, title_str, desc_str) in enumerate(steps):
        sy = t + (i * step_h) + 0.15

        # Circle Badge
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l + 0.25), Inches(sy), Inches(0.42), Inches(0.42))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_ORANGE
        circle.line.fill.background()
        tf_c = circle.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.text = str(num_str)
        p_c.alignment = PP_ALIGN.CENTER
        p_c.font.name = FONT_BODY
        p_c.font.size = Pt(12)
        p_c.font.bold = True
        p_c.font.color.rgb = COLOR_WHITE

        # Step Text
        tb = slide.shapes.add_textbox(Inches(l + 0.78), Inches(sy - 0.05), Inches(w - 0.95), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title_str
        p1.font.name = FONT_BODY
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_DARK

        p2 = tf.add_paragraph()
        p2.text = desc_str
        p2.font.name = FONT_BODY
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = COLOR_MUTED


# ── Presentation Builder ──────────────────────────────────────────────────────

def build_presentation():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ==========================================================================
    # SLIDE 1 : COVER SLIDE (Matching "The Caillou Project" Slide 1)
    # ==========================================================================
    s1 = prs.slides.add_slide(blank_layout)

    # Full-bleed Background Image
    if BG_TITLE.exists():
        s1.shapes.add_picture(str(BG_TITLE), 0, 0, width=Inches(13.333), height=Inches(7.5))
    else:
        bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg1.fill.solid()
        bg1.fill.fore_color.rgb = RGBColor(40, 50, 60)

    # Badges / Logos on Cover
    if ESTACA_PATH.exists():
        s1.shapes.add_picture(str(ESTACA_PATH), Inches(0.6), Inches(0.5), height=Inches(1.2))
    if OWF_PATH.exists():
        s1.shapes.add_picture(str(OWF_PATH), Inches(11.2), Inches(0.5), height=Inches(1.2))

    # White Horizontal Banner Across Slide
    banner_y = 5.2
    banner_h = 1.1
    banner = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(banner_y), Inches(13.333), Inches(banner_h))
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLOR_WHITE
    banner.line.color.rgb = RGBColor(200, 200, 200)
    banner.line.width = Pt(1)

    # Title Text Inside Banner
    tb_t = s1.shapes.add_textbox(0, Inches(banner_y + 0.18), Inches(13.333), Inches(0.75))
    p_t = tb_t.text_frame.paragraphs[0]
    p_t.text = "The JARVIS / AOUDA Project"
    p_t.alignment = PP_ALIGN.CENTER
    p_t.font.name = FONT_HEADING
    p_t.font.size = Pt(36)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_DARK

    # Subtitle Text Below Banner
    tb_sub = s1.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4))
    p_sub = tb_sub.text_frame.paragraphs[0]
    p_sub.text = "ESTACA x ÖWF AMADEE-27 Analog Mars Mission  |  Internship Project 2026"
    p_sub.alignment = PP_ALIGN.RIGHT
    p_sub.font.name = FONT_BODY
    p_sub.font.size = Pt(13)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_WHITE

    # ==========================================================================
    # SLIDE 2 : CONTEXT (Matching Slide 2 of Reference PDF)
    # ==========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "Context")

    # Body Text Paragraph 1
    p1_text = (
        "The upcoming AMADEE-27 mission, organized by the Austrian Space Forum (OeWF), "
        "represents a major milestone in high-fidelity Mars analog simulations. Based in Portugal, "
        "the mission will test key instruments and exploration workflows critical to human planetary exploration. "
        "During Extravehicular Activities (EVAs), analog astronauts wear the complex AOUDA spacesuit simulator. "
        "However, operating technical instruments and consulting paper manuals while wearing heavy pressurized gloves "
        "and restricted visors introduces severe operational friction."
    )
    add_text_paragraph(s2, p1_text, 1.0, 1.6, 11.333, 2.2, font_name=FONT_BODY, font_size=15, space_after=Pt(16))

    # Body Text Paragraph 2
    p2_text = (
        "To eliminate screen dependency and manual handling risks, this project introduces 'JARVIS / AOUDA', "
        "a fully autonomous, 100% offline voice assistant specifically tailored for Mars EVA procedures. "
        "By providing deterministic, zero-hallucination voice guidance, session memory for scientific payloads "
        "such as the 'Caillou' Raman Spectrometer, and real-time telemetry monitoring, JARVIS optimizes astronaut "
        "safety, procedural efficiency, and scientific data collection under extreme operational constraints."
    )
    add_text_paragraph(s2, p2_text, 1.0, 4.0, 11.333, 2.5, font_name=FONT_BODY, font_size=15, space_after=Pt(14))

    # ==========================================================================
    # SLIDE 3 : CONTENT / AGENDA (Matching Slide 3 of Reference PDF)
    # ==========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "Content")

    agenda_items = [
        ("System Requirements & Operational Scenario", "Functional constraints, EVA workflow, and telemetry integration."),
        ("Dual-Brain Software Architecture", "Deterministic Little Brain reflex core & Big Brain procedure engine."),
        ("Voice Processing Pipeline & Offline STT/TTS", "Faster-Whisper transcription, phonetic correction, and Piper TTS."),
        ("Mission Scenario & Field Demo ('Caillou')", "Step-by-step interactive demonstration with 'Caillou' spectrometer payload."),
        ("Hardware Setup & System Validation", "Raspberry Pi SBC core, power autonomy, and 33/33 unit test suite results.")
    ]

    for idx, (title, sub) in enumerate(agenda_items):
        y_pos = 1.6 + (idx * 1.05)

        # Icon / Symbol
        tb_ic = s3.shapes.add_textbox(Inches(1.2), Inches(y_pos), Inches(0.6), Inches(0.5))
        p_ic = tb_ic.text_frame.paragraphs[0]
        p_ic.text = "⚙"
        p_ic.font.name = FONT_BODY
        p_ic.font.size = Pt(20)
        p_ic.font.color.rgb = COLOR_DARK

        # Agenda Text
        tb_a = s3.shapes.add_textbox(Inches(1.8), Inches(y_pos), Inches(10.5), Inches(0.8))
        tf_a = tb_a.text_frame
        tf_a.word_wrap = True

        p_title = tf_a.paragraphs[0]
        p_title.text = title
        p_title.font.name = FONT_HEADING
        p_title.font.size = Pt(17)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_DARK

        p_sub = tf_a.add_paragraph()
        p_sub.text = sub
        p_sub.font.name = FONT_BODY
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = COLOR_MUTED

    # ==========================================================================
    # SLIDE 4 : SYSTEM REQUIREMENTS (Matching Slide 4 of Reference PDF)
    # ==========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "System Requirements & Operational Scenario")

    # Left Subtitle
    add_text_paragraph(s4, "Functional & Data Workflow", 0.8, 1.4, 7.5, 0.5,
                       font_name=FONT_HEADING, font_size=20, bold=True)

    # Left Bullets
    bullets_s4 = [
        ("Realistic Workflow", "Replicate the complete setup and voice data path of an EVA mission assistant without cloud connectivity."),
        ("Targeting & Context", "Integrate voice-driven procedure navigation ('Caillou' spectrometer) with hands-free execution."),
        ("Synthetic Output", "Generate deterministic step responses from local libraries and stream real-time suit telemetry to the astronaut.")
    ]
    add_bullet_list(s4, bullets_s4, 0.8, 2.3, 7.5, 4.5, font_size=14, space_after=Pt(22))

    # Right Sidebar Workflow
    steps_s4 = [
        (1, "Arrival & Wake-Word", "Astronaute initiates assistant via wake phrase 'AOUDA'."),
        (2, "STT Speech Capture", "Local Faster-Whisper converts streaming audio to raw text."),
        (3, "Phonetic Correction", "Regex engine fixes misheard payload words (e.g. 'call you' -> 'caillou')."),
        (4, "Dual-Brain Routing", "Little Brain verifies vital suit telemetry before Big Brain checks procedures."),
        (5, "Procedure Execution", "System retrieves exact procedure step & maintains session memory."),
        (6, "Audio Synthesis", "Piper TTS generates high-clarity voice output in suit headset.")
    ]
    add_process_sidebar(s4, steps_s4, l=8.7, t=1.4, w=4.1, h=5.6)

    # ==========================================================================
    # SLIDE 5 : SOFTWARE ARCHITECTURE (Matching Slide 5 Layout)
    # ==========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "Software & Dual-Brain Architecture")

    # Left Subtitle
    add_text_paragraph(s5, "Dual-Brain Core Engine", 0.8, 1.4, 7.5, 0.5,
                       font_name=FONT_HEADING, font_size=20, bold=True)

    # Left Bullets
    bullets_s5 = [
        ("Little Brain (Reflex Engine)", "Handles priority telemetry checks (O2, CO2, heart rate, battery) and emergency alerts with <1ms execution latencies."),
        ("Big Brain (Procedure Engine)", "Ingests local mission documents (Markdown, YAML, JSON) and provides step-by-step navigation (next, back, repeat)."),
        ("Deterministic Logic", "Strict rule-based evaluation guarantees zero hallucination and 100% predictable responses during mission-critical tasks.")
    ]
    add_bullet_list(s5, bullets_s5, 0.8, 2.3, 7.5, 4.5, font_size=14, space_after=Pt(22))

    # Right Process Steps
    steps_s5 = [
        (1, "Microphone Input", "Raw 16kHz audio stream captured inside AOUDA helmet."),
        (2, "Reflex Filter", "Little Brain inspects stream for instant emergency triggers."),
        (3, "Procedure Context", "Big Brain identifies active payload ('Caillou') from session memory."),
        (4, "Local Document Search", "Queries structured Markdown files without external network calls."),
        (5, "Format Response", "Assembles exact step instructions and telemetry state."),
        (6, "Speaker Output", "Piper neural TTS renders audio stream to astronaut headset.")
    ]
    add_process_sidebar(s5, steps_s5, l=8.7, t=1.4, w=4.1, h=5.6)

    # ==========================================================================
    # SLIDE 6 : MISSION SCENARIO & DEMO (Matching Slide 7 of PDF)
    # ==========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "Mission Scenario : 'Caillou' Spectrometer Demo")

    # 4 Scenario Step Cards arranged in 2x2 Grid
    grid_coords = [
        (0.8, 1.5, "Step 1: Payload Selection", "Astronaut: \"AOUDA, Caillou step 4\"\nJARVIS: \"Step 4 — Laser Targeting: Activate laser beam via remote interface. Duration: 10s.\""),
        (6.8, 1.5, "Step 2: Contextual Follow-up", "Astronaut: \"What do I do after that?\"\nJARVIS: \"Step 5 — Alignment: Align laser dot with target rock sample surface.\""),
        (0.8, 4.2, "Step 3: Repetition & Clarification", "Astronaut: \"Repeat please\"\nJARVIS: \"Step 5 — Alignment: Align laser dot with target rock sample surface.\""),
        (6.8, 4.2, "Step 4: Live Telemetry Check", "Astronaut: \"What is my oxygen level?\"\nJARVIS: \"Oxygen: 20.9% — Normal. Heart Rate: 85 bpm — Normal.\"")
    ]

    for x, y, title, dialogue in grid_coords:
        # Card outline box
        card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.7), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BG_CARD
        card.line.color.rgb = COLOR_BORDER
        card.line.width = Pt(1)

        # Card Title Banner
        tb_ct = s6.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(5.3), Inches(0.4))
        p_ct = tb_ct.text_frame.paragraphs[0]
        p_ct.text = title
        p_ct.font.name = FONT_HEADING
        p_ct.font.size = Pt(14)
        p_ct.font.bold = True
        p_ct.font.color.rgb = COLOR_DARK

        # Dialogue text
        tb_cd = s6.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.55), Inches(5.3), Inches(1.6))
        tf_cd = tb_cd.text_frame
        tf_cd.word_wrap = True

        lines = dialogue.split('\n')
        for i, line in enumerate(lines):
            p = tf_cd.paragraphs[0] if i == 0 else tf_cd.add_paragraph()
            p.text = line
            p.font.name = FONT_BODY
            p.font.size = Pt(12)
            p.space_after = Pt(4)
            if "Astronaut:" in line:
                p.font.bold = True
                p.font.color.rgb = COLOR_BLUE_ACC
            else:
                p.font.color.rgb = COLOR_DARK

    # ==========================================================================
    # SLIDE 7 : HARDWARE & MECHANICAL DESIGN (Matching Slide 8 of Reference PDF)
    # ==========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Hardware & Mechanical Design")

    hw_cards = [
        (0.8, 1.6, "SBC Core Unit", [
            ("Processor", "Raspberry Pi 4 / Embedded Linux"),
            ("STT Engine", "CTranslate2 Faster-Whisper (Int8)"),
            ("Audio Interface", "16kHz USB Microphone & Headset"),
            ("Offline Storage", "Local SSD with preloaded procedures")
        ]),
        (4.8, 1.6, "Power & Autonomy", [
            ("Power Source", "Lightweight commercial power bank"),
            ("Runtime", "2+ hours continuous EVA operation"),
            ("Thermal Shield", "Custom insulation for field durability"),
            ("Monitoring", "Voltage telemetry fed to Little Brain")
        ]),
        (8.8, 1.6, "Ergonomics & Integration", [
            ("Suit Mounting", "Hands-free helmet headset connection"),
            ("EVA Usability", "Zero touch required during operations"),
            ("Field Durability", "Dust-resistant sealed enclosure"),
            ("Mission Testing", "Validated for ÖWF AMADEE analog suit")
        ])
    ]

    for x, y, title, specs in hw_cards:
        card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.7), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BG_CARD
        card.line.color.rgb = COLOR_BORDER
        card.line.width = Pt(1)

        # Header box inside card
        hb = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(3.7), Inches(0.6))
        hb.fill.solid()
        hb.fill.fore_color.rgb = COLOR_DARK
        hb.line.fill.background()

        tb_h = s7.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.1), Inches(3.4), Inches(0.4))
        p_h = tb_h.text_frame.paragraphs[0]
        p_h.text = title
        p_h.font.name = FONT_HEADING
        p_h.font.size = Pt(14)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE

        # Specification list
        tb_sp = s7.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.8), Inches(3.3), Inches(4.2))
        tf_sp = tb_sp.text_frame
        tf_sp.word_wrap = True

        for idx, (label, val) in enumerate(specs):
            p = tf_sp.paragraphs[0] if idx == 0 else tf_sp.add_paragraph()
            p.space_after = Pt(14)

            r1 = p.add_run()
            r1.text = f"•  {label}:\n"
            r1.font.name = FONT_BODY
            r1.font.size = Pt(12)
            r1.font.bold = True
            r1.font.color.rgb = COLOR_DARK

            r2 = p.add_run()
            r2.text = val
            r2.font.name = FONT_BODY
            r2.font.size = Pt(11.5)
            r2.font.color.rgb = COLOR_MUTED

    # ==========================================================================
    # SLIDE 8 : CONCLUSION & RESULTS
    # ==========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "Conclusion & Validation Results")

    # 3 Large Metric Cards
    metrics = [
        ("33 / 33", "PyTest Unit Tests Passed", COLOR_BLUE_ACC),
        ("< 1.5 s", "End-to-End Voice Latency", COLOR_DARK),
        ("0 %", "Hallucination Rate (100% Deterministic)", COLOR_RED_ACC)
    ]

    for i, (val, label, color) in enumerate(metrics):
        x = 0.8 + (i * 4.0)
        card = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(3.7), Inches(2.0))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BG_CARD
        card.line.color.rgb = COLOR_BORDER
        card.line.width = Pt(1)

        tb_v = s8.shapes.add_textbox(Inches(x + 0.1), Inches(1.65), Inches(3.5), Inches(0.9))
        p_v = tb_v.text_frame.paragraphs[0]
        p_v.text = val
        p_v.alignment = PP_ALIGN.CENTER
        p_v.font.name = FONT_HEADING
        p_v.font.size = Pt(36)
        p_v.font.bold = True
        p_v.font.color.rgb = color

        tb_l = s8.shapes.add_textbox(Inches(x + 0.1), Inches(2.6), Inches(3.5), Inches(0.7))
        p_l = tb_l.text_frame.paragraphs[0]
        p_l.text = label
        p_l.alignment = PP_ALIGN.CENTER
        p_l.font.name = FONT_BODY
        p_l.font.size = Pt(12)
        p_l.font.bold = True
        p_l.font.color.rgb = COLOR_DARK

    # Summary Box
    card_sum = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.8), Inches(11.7), Inches(3.1))
    card_sum.fill.solid()
    card_sum.fill.fore_color.rgb = COLOR_BG_CARD
    card_sum.line.color.rgb = COLOR_BORDER
    card_sum.line.width = Pt(1)

    tb_st = s8.shapes.add_textbox(Inches(1.1), Inches(4.0), Inches(11.1), Inches(0.4))
    p_st = tb_st.text_frame.paragraphs[0]
    p_st.text = "Summary & AMADEE-27 Perspectives"
    p_st.font.name = FONT_HEADING
    p_st.font.size = Pt(16)
    p_st.font.bold = True
    p_st.font.color.rgb = COLOR_DARK

    bullets_s8 = [
        ("Proven Reliability", "Fully local processing ensures zero dependency on external network or cloud services."),
        ("Mission Expansion", "Adding new procedures requires simply dropping Markdown or YAML files into the knowledge base."),
        ("Next Steps", "Integrate with suit HUD visual display and conduct field trials during ÖWF AMADEE Mars analog simulations.")
    ]
    add_bullet_list(s8, bullets_s8, 1.1, 4.5, 11.1, 2.2, font_size=13.5, space_after=Pt(10))

    # Save output
    out_path = Path("presentation_projet_jarvis_aouda.pptx")
    prs.save(str(out_path))
    print(f"[OK] Presentation regenerated successfully: {out_path.resolve()}")


if __name__ == "__main__":
    build_presentation()
